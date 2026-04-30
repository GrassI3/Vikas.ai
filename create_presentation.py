"""
Vikas.ai — Hackathon Presentation Generator
Creates a professional 5-slide PPTX for the hackathon demo.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Constants ────────────────────────────────────────────────
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Colors
BG_DARK       = RGBColor(0x0F, 0x0F, 0x1A)    # Very dark navy
BG_CARD       = RGBColor(0x1A, 0x1A, 0x2E)    # Card dark
ACCENT_INDIGO = RGBColor(0x63, 0x66, 0xF1)    # Primary indigo
ACCENT_CYAN   = RGBColor(0x06, 0xB6, 0xD4)    # Cyan accent
ACCENT_EMERALD= RGBColor(0x10, 0xB9, 0x81)    # Emerald
ACCENT_AMBER  = RGBColor(0xF5, 0x9E, 0x0B)    # Amber
ACCENT_ROSE   = RGBColor(0xF4, 0x3F, 0x5E)    # Rose
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_PRIMARY  = RGBColor(0xE2, 0xE8, 0xF0)
TEXT_MUTED    = RGBColor(0x94, 0xA3, 0xB8)
TEXT_DIM      = RGBColor(0x64, 0x74, 0x8B)


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    # Adjust corner rounding
    shape.adjustments[0] = 0.05
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=TEXT_PRIMARY, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multi_text(slide, left, top, width, height, lines, default_size=16,
                   default_color=TEXT_PRIMARY, line_spacing=1.4):
    """Add a text box with multiple styled paragraphs.
    lines: list of dicts with keys: text, size, color, bold, spacing_before (optional)
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line_data in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = line_data.get("text", "")
        p.font.size = Pt(line_data.get("size", default_size))
        p.font.color.rgb = line_data.get("color", default_color)
        p.font.bold = line_data.get("bold", False)
        p.font.name = line_data.get("font", "Calibri")
        p.alignment = line_data.get("align", PP_ALIGN.LEFT)
        if "spacing_before" in line_data:
            p.space_before = Pt(line_data["spacing_before"])
        if "spacing_after" in line_data:
            p.space_after = Pt(line_data["spacing_after"])
    return txBox


def add_accent_line(slide, left, top, width, color=ACCENT_INDIGO):
    """Add a horizontal accent line."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(4)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — Title Slide
# ═══════════════════════════════════════════════════════════════
def create_slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_bg(slide, BG_DARK)

    # Accent line at top
    add_accent_line(slide, Inches(0), Inches(0), SLIDE_WIDTH, ACCENT_INDIGO)

    # Category badge
    add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(0.5),
                 "10  HEALTH · SOCIAL IMPACT", font_size=14, color=ACCENT_CYAN,
                 bold=True, font_name="Calibri")

    # Main title
    add_text_box(slide, Inches(1), Inches(1.9), Inches(11), Inches(1.2),
                 "Vikas.ai", font_size=60, color=WHITE, bold=True,
                 font_name="Calibri")

    # Subtitle
    add_multi_text(slide, Inches(1), Inches(3.2), Inches(10), Inches(1.5), [
        {"text": "Explainable, Voice-Driven Decision Support for", "size": 24,
         "color": TEXT_PRIMARY, "bold": False},
        {"text": "Underserved Communities in India", "size": 24,
         "color": ACCENT_CYAN, "bold": True, "spacing_before": 4},
    ])

    # Divider
    add_accent_line(slide, Inches(1), Inches(4.5), Inches(3), ACCENT_INDIGO)

    # Key highlights in a row
    highlights = [
        ("📞", "Voice-First", "No app, no internet"),
        ("🌐", "Multilingual", "8+ Indian languages"),
        ("🧠", "Explainable AI", "Cited medical reasoning"),
        ("🛡️", "Safety-First", "Emergency protocols"),
    ]

    for i, (icon, title, sub) in enumerate(highlights):
        x = Inches(1 + i * 2.9)
        y = Inches(5.0)
        card = add_shape(slide, x, y, Inches(2.6), Inches(1.5), BG_CARD,
                         RGBColor(0x2D, 0x2D, 0x4A))
        add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(2.2), Inches(0.4),
                     icon, font_size=28, alignment=PP_ALIGN.LEFT)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.55), Inches(2.2), Inches(0.35),
                     title, font_size=15, color=WHITE, bold=True)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.9), Inches(2.2), Inches(0.4),
                     sub, font_size=11, color=TEXT_MUTED)

    # Team / Footer
    add_text_box(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.4),
                 "Hackathon 2026  ·  PS #10  ·  Health & Social Impact Track",
                 font_size=12, color=TEXT_DIM, alignment=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — Problem Statement
# ═══════════════════════════════════════════════════════════════
def create_slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_accent_line(slide, Inches(0), Inches(0), SLIDE_WIDTH, ACCENT_ROSE)

    # Section tag
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(4), Inches(0.4),
                 "01  PROBLEM STATEMENT", font_size=13, color=ACCENT_ROSE, bold=True)

    # Title
    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8),
                 "People Don't Need More Information —", font_size=36,
                 color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.8),
                 "They Need Clear, Contextual Decisions", font_size=36,
                 color=ACCENT_CYAN, bold=True)

    # Problem cards — left column
    problems = [
        ("🚫", "Access Barriers",
         "900M+ Indians lack reliable internet. Existing health apps require smartphones, data plans, and literacy — excluding the most vulnerable."),
        ("🤯", "Information Overload",
         "Google returns 500M results for 'headache'. Underserved users cannot filter peer-reviewed evidence from misinformation."),
        ("⏱️", "Critical Delay",
         "Delayed medical triage costs lives. Rural users travel 50+ km to discover their condition was manageable at home — or miss emergencies."),
    ]

    for i, (icon, title, desc) in enumerate(problems):
        x = Inches(0.8)
        y = Inches(2.6 + i * 1.45)
        card = add_shape(slide, x, y, Inches(6.2), Inches(1.25), BG_CARD,
                         RGBColor(0x2D, 0x2D, 0x4A))
        add_text_box(slide, x + Inches(0.15), y + Inches(0.08), Inches(0.5), Inches(0.4),
                     icon, font_size=22)
        add_text_box(slide, x + Inches(0.6), y + Inches(0.08), Inches(5.3), Inches(0.35),
                     title, font_size=16, color=WHITE, bold=True)
        add_text_box(slide, x + Inches(0.6), y + Inches(0.45), Inches(5.3), Inches(0.7),
                     desc, font_size=12, color=TEXT_MUTED)

    # Right column — Impact stats
    stats = [
        ("900M+", "Indians without\nreliable internet", ACCENT_ROSE),
        ("60%", "Rely on unverified\nhealth advice", ACCENT_AMBER),
        ("3.5 Hrs", "Average rural travel\nfor basic consult", ACCENT_INDIGO),
        ("22", "Official Indian\nlanguages", ACCENT_CYAN),
    ]

    for i, (value, label, color) in enumerate(stats):
        x = Inches(7.4)
        y = Inches(2.6 + i * 1.1)
        card = add_shape(slide, x, y, Inches(5.2), Inches(0.9), BG_CARD,
                         RGBColor(0x2D, 0x2D, 0x4A))
        add_text_box(slide, x + Inches(0.3), y + Inches(0.08), Inches(2), Inches(0.5),
                     value, font_size=28, color=color, bold=True)
        add_text_box(slide, x + Inches(2.4), y + Inches(0.15), Inches(2.5), Inches(0.6),
                     label.replace("\n", " "), font_size=12, color=TEXT_MUTED)

    # Bottom quote
    add_text_box(slide, Inches(0.8), Inches(6.7), Inches(11.5), Inches(0.5),
                 "\"The most vulnerable people in the world can be reached with just a phone call — "
                 "no app, no internet, no literacy required.\"",
                 font_size=13, color=TEXT_DIM, alignment=PP_ALIGN.CENTER,
                 font_name="Calibri")


# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — Solution
# ═══════════════════════════════════════════════════════════════
def create_slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_accent_line(slide, Inches(0), Inches(0), SLIDE_WIDTH, ACCENT_EMERALD)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(4), Inches(0.4),
                 "02  OUR SOLUTION", font_size=13, color=ACCENT_EMERALD, bold=True)

    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8),
                 "Vikas.ai — Just Dial & Talk", font_size=38,
                 color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.7), Inches(10), Inches(0.5),
                 "An AI-powered telephonic assistant providing explainable, evidence-based health & civic guidance — "
                 "accessible to anyone with a basic phone.",
                 font_size=15, color=TEXT_MUTED)

    # How it works — 4 step flow
    steps = [
        ("📞", "Call the Number", "Dial the helpline from any phone — feature phone, smartphone, landline. No app or internet needed.", ACCENT_INDIGO),
        ("🗣️", "Speak Naturally", "Talk in Hindi, Marathi, Tamil, Telugu, English, or 8+ languages. AI auto-detects your language.", ACCENT_CYAN),
        ("🧠", "AI Reasons & Responds", "Multi-agent pipeline performs clinical triage, retrieves PubMed evidence, and builds an explainable response.", ACCENT_EMERALD),
        ("✅", "Get Actionable Guidance", "Receive clear next-steps: self-care routines, when to see a doctor, emergency escalation with helpline numbers.", ACCENT_AMBER),
    ]

    for i, (icon, title, desc, color) in enumerate(steps):
        x = Inches(0.5 + i * 3.1)
        y = Inches(2.6)
        card = add_shape(slide, x, y, Inches(2.9), Inches(2.4), BG_CARD, color)
        # Step number
        add_text_box(slide, x + Inches(0.15), y + Inches(0.1), Inches(2.6), Inches(0.35),
                     f"{icon}  Step {i+1}", font_size=14, color=color, bold=True)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.5), Inches(2.6), Inches(0.4),
                     title, font_size=17, color=WHITE, bold=True)
        add_text_box(slide, x + Inches(0.15), y + Inches(1.0), Inches(2.6), Inches(1.2),
                     desc, font_size=12, color=TEXT_MUTED)

    # Key differentiators
    diffs = [
        ("Explainable Reasoning", "Full chain-of-thought shown — not a black box. Users understand WHY."),
        ("Safety Guardrails", "Emergency keyword detection, hard-coded disclaimers, prescription blocking."),
        ("Live Intervention", "Guided breathing exercises, panic de-escalation, and suicide hotline routing."),
        ("SMS Fallback", "If the call drops, the response is delivered via SMS automatically."),
    ]

    for i, (title, desc) in enumerate(diffs):
        x = Inches(0.5 + i * 3.1)
        y = Inches(5.3)
        card = add_shape(slide, x, y, Inches(2.9), Inches(1.2), BG_CARD,
                         RGBColor(0x2D, 0x2D, 0x4A))
        add_text_box(slide, x + Inches(0.15), y + Inches(0.1), Inches(2.6), Inches(0.3),
                     f"✦ {title}", font_size=13, color=ACCENT_CYAN, bold=True)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.45), Inches(2.6), Inches(0.6),
                     desc, font_size=11, color=TEXT_MUTED)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — Tech Stack & Architecture
# ═══════════════════════════════════════════════════════════════
def create_slide_4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_accent_line(slide, Inches(0), Inches(0), SLIDE_WIDTH, ACCENT_INDIGO)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.4),
                 "03  TECH STACK & ARCHITECTURE", font_size=13, color=ACCENT_INDIGO, bold=True)

    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8),
                 "End-to-End Voice AI Pipeline", font_size=36,
                 color=WHITE, bold=True)

    # Architecture flow — horizontal pipeline
    arch_nodes = [
        ("📞", "User Phone\n(PSTN/VoIP)", "Any phone\nin India", ACCENT_AMBER),
        ("🎙️", "Vapi Gateway\n+ Twilio", "Voice ↔ Text\nOrchestration", ACCENT_INDIGO),
        ("📝", "Deepgram\nNova 3", "Multilingual\nTranscription", ACCENT_CYAN),
        ("🧠", "Groq + LLama\n3.1 8B Instant", "LLM Reasoning\n& Triage", ACCENT_EMERALD),
        ("🔊", "ElevenLabs\nVikram", "Multilingual\nV2 TTS", ACCENT_ROSE),
    ]

    for i, (icon, title, sub, color) in enumerate(arch_nodes):
        x = Inches(0.3 + i * 2.55)
        y = Inches(1.9)
        card = add_shape(slide, x, y, Inches(2.3), Inches(1.4), BG_CARD, color)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.05), Inches(2.0), Inches(0.4),
                     icon, font_size=24, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.05), y + Inches(0.45), Inches(2.2), Inches(0.5),
                     title.replace("\n", " "), font_size=13, color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.05), y + Inches(0.9), Inches(2.2), Inches(0.4),
                     sub.replace("\n", " "), font_size=10, color=TEXT_MUTED,
                     alignment=PP_ALIGN.CENTER)

    # Arrow connectors between nodes (simple text arrows)
    for i in range(4):
        x = Inches(2.45 + i * 2.55)
        add_text_box(slide, x, Inches(2.3), Inches(0.4), Inches(0.4),
                     "→", font_size=24, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

    # Tech stack detail cards — bottom grid
    stack_sections = [
        ("Telephony Layer", [
            "Vapi AI — Voice assistant orchestration",
            "Twilio — Phone numbers & SMS OTP auth",
            "ngrok — Secure webhook tunneling",
        ], ACCENT_INDIGO),
        ("AI / ML Layer", [
            "Groq — Ultra-fast LLM inference (LLama 3.1 8B)",
            "LangGraph — Multi-agent state machine",
            "ChromaDB — Vector RAG retrieval",
            "Sentence Transformers — Embeddings",
        ], ACCENT_EMERALD),
        ("Voice / NLP Layer", [
            "Deepgram Nova 3 — Multilingual ASR",
            "ElevenLabs Multilingual V2 — TTS (Vikram)",
            "8+ Indian languages supported",
        ], ACCENT_CYAN),
        ("Backend / Infra", [
            "FastAPI — Async Python backend",
            "PubMed API — Live medical evidence",
            "Safety Guardrails — Emergency protocols",
            "SMS Fallback — Twilio delivery",
        ], ACCENT_AMBER),
    ]

    for i, (title, items, color) in enumerate(stack_sections):
        x = Inches(0.3 + i * 3.2)
        y = Inches(3.7)
        card = add_shape(slide, x, y, Inches(3.0), Inches(3.3), BG_CARD,
                         RGBColor(0x2D, 0x2D, 0x4A))
        add_text_box(slide, x + Inches(0.15), y + Inches(0.1), Inches(2.7), Inches(0.35),
                     title, font_size=15, color=color, bold=True)

        # Accent line under title
        add_shape(slide, x + Inches(0.15), y + Inches(0.45), Inches(2.0), Pt(2), color)

        for j, item in enumerate(items):
            add_text_box(slide, x + Inches(0.15), y + Inches(0.6 + j * 0.55),
                         Inches(2.7), Inches(0.5),
                         f"▸ {item}", font_size=11, color=TEXT_MUTED)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — Demo & Thank You
# ═══════════════════════════════════════════════════════════════
def create_slide_5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_accent_line(slide, Inches(0), Inches(0), SLIDE_WIDTH, ACCENT_CYAN)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(4), Inches(0.4),
                 "04  LIVE DEMO & WRAP-UP", font_size=13, color=ACCENT_CYAN, bold=True)

    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8),
                 "Try It Live — Call Vikas Now", font_size=38,
                 color=WHITE, bold=True)

    # Phone numbers
    numbers = [
        ("+1 (631) 490-9141", "Vapi Primary · SMS Auth", ACCENT_INDIGO),
        ("+1 (938) 902-2543", "Twilio · SMS Auth", ACCENT_CYAN),
        ("+1 (516) 667-0818", "Vapi Backup", ACCENT_EMERALD),
        ("+98993 56642", "Vapi India", ACCENT_AMBER),
    ]

    for i, (number, label, color) in enumerate(numbers):
        x = Inches(0.8 + i * 3.1)
        y = Inches(2.0)
        card = add_shape(slide, x, y, Inches(2.8), Inches(1.2), BG_CARD, color)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.15), Inches(2.5), Inches(0.5),
                     number, font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.7), Inches(2.5), Inches(0.35),
                     label, font_size=11, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

    # Demo script
    add_text_box(slide, Inches(0.8), Inches(3.5), Inches(6), Inches(0.4),
                 "Live Demo Script", font_size=18, color=ACCENT_CYAN, bold=True)

    demo_lines = [
        "1.  Call the number — Vikas greets in Hindi",
        "2.  Say symptoms in Hindi, Marathi, or English",
        "3.  AI reasons through PubMed evidence in real-time",
        "4.  Get actionable guidance with cited sources",
        "5.  Test emergency response — \"mujhe saans nahi aa rahi\"",
    ]

    for i, line in enumerate(demo_lines):
        add_text_box(slide, Inches(1.0), Inches(4.0 + i * 0.4), Inches(5.5), Inches(0.35),
                     line, font_size=14, color=TEXT_PRIMARY)

    # PS Requirements checklist
    add_text_box(slide, Inches(7.2), Inches(3.5), Inches(5.5), Inches(0.4),
                 "PS Requirements — All Covered", font_size=18, color=ACCENT_EMERALD, bold=True)

    reqs = [
        ("✅", "Conversational input flow", "Voice-first telephony"),
        ("✅", "Reasoning engine", "LangGraph multi-agent CoT"),
        ("✅", "Safety disclaimers", "Emergency protocols + guardrails"),
        ("✅", "Explainable output", "Full reasoning chain visible"),
        ("✅", "Multilingual support", "8+ Indian languages"),
        ("✅", "Voice input & output", "Deepgram ASR + ElevenLabs TTS"),
    ]

    for i, (check, req, detail) in enumerate(reqs):
        y = Inches(4.0 + i * 0.4)
        add_text_box(slide, Inches(7.2), y, Inches(0.4), Inches(0.35),
                     check, font_size=14)
        add_text_box(slide, Inches(7.7), y, Inches(2.2), Inches(0.35),
                     req, font_size=13, color=WHITE, bold=True)
        add_text_box(slide, Inches(10.0), y, Inches(2.8), Inches(0.35),
                     detail, font_size=11, color=TEXT_MUTED)

    # Thank you
    add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.6),
                 "Thank You — Questions?",
                 font_size=28, color=ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════
def main():
    prs = Presentation()

    # Set slide dimensions to widescreen 16:9
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    create_slide_1(prs)
    create_slide_2(prs)
    create_slide_3(prs)
    create_slide_4(prs)
    create_slide_5(prs)

    output_path = "Vikas_AI_Presentation.pptx"
    prs.save(output_path)
    print(f"✅ Presentation saved to: {output_path}")
    print(f"   Slides: {len(prs.slides)}")
    print(f"   Format: 16:9 Widescreen")


if __name__ == "__main__":
    main()
